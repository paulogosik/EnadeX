"""
Confere que os números de benchmark citados no documento (.docx) e na
apresentação (.pptx) existem no banco — nas views `v_benchmark_resumo` e
`v_benchmark_metricas` — com tolerância de arredondamento. É o aceite (e) da
Fase 1: nenhum número sem linha correspondente na view.

O que é verificado
  * Só tokens com unidade de benchmark: segundos ("112,93 s"), speedup
    ("2,4632×" / "1,52×"), eficiência/percentual ("61,58 %") e throughput
    ("116.487 linhas/s"). Outros números do texto (contagens, anos, versões)
    ficam fora do escopo.
  * Lista negra: os valores da rodada perdida (243,55 / 198,57 / 218,97 s,
    1,2265× / 1,227×, 61,3 % / 27,8 %) e o card errado do documento de 21/08
    (2,6157×, 65,39 %) NÃO podem aparecer. Falha imediata.
  * Rastreabilidade: cada token precisa bater (2 casas para s e %, 4 ou 2 para
    ×) com algum valor das views — da campanha oficial OU de qualquer
    execução histórica (a seção de histórico cita 21/06 e 21/08 de propósito)
    — ou com um teto de escalonamento do JSON de `analise_escalonamento.py`.
    Tokens não rastreáveis são listados; com --estrito, também falham.

Uso
  python docs/geradores/verificar_numeros.py --pptx docs/geradores/out/apresentacao.pptx \
      --docx docs/geradores/out/documento.docx [--escalonamento out/escalonamento.json] [--estrito]
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

RAIZ = Path(__file__).resolve().parents[3]  # raiz do EnadeX
sys.path.insert(0, str(RAIZ))

from enade_time.etl.bench_db import conectar  # noqa: E402

LISTA_NEGRA = {
    "s": {"243,55", "198,57", "218,97", "243,5", "198,6"},
    "x": {"1,2265", "1,227", "1,1122", "1,112"},
    "%": {"61,32", "61,3", "27,81", "27,8"},
    "lps": {"21.927", "26.893", "24.388"},
}

RE_SEG = re.compile(r"(?<![\d.,])(\d{1,3}(?:\.\d{3})*,\d{1,4})\s?s\b")
RE_X = re.compile(r"(?<![\d.,])(\d{1,2},\d{1,4})\s?[×x]\b")
RE_PCT = re.compile(r"(?<![\d.,])(\d{1,3},\d{1,2})\s?%")
RE_LPS = re.compile(r"(?<![\d.,])(\d{1,3}(?:\.\d{3})+|\d{4,6})\s?linhas/s")


def fmt_br(v: float, casas: int) -> str:
    s = f"{v:,.{casas}f}"
    return s.replace(",", "X").replace(".", ",").replace("X", ".")


def texto_pptx(caminho: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(caminho))
    partes: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                partes.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    partes.append(" | ".join(c.text for c in row.cells))
            if getattr(shape, "has_chart", False) and shape.has_chart:
                for plot in shape.chart.plots:
                    for serie in plot.series:
                        partes.append(" ".join(f"{v:.2f}".replace(".", ",") + " s" for v in serie.values if v is not None))
        if slide.has_notes_slide:
            partes.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(partes)


def texto_docx(caminho: Path) -> str:
    from docx import Document
    doc = Document(str(caminho))
    partes = [p.text for p in doc.paragraphs]
    for t in doc.tables:
        for row in t.rows:
            partes.append(" | ".join(c.text for c in row.cells))
    return "\n".join(partes)


def valores_permitidos(conn, escalonamento: dict | None) -> dict[str, set[str]]:
    seg: set[str] = set(); x: set[str] = set(); pct: set[str] = set(); lps: set[str] = set()

    def add_tempo(v):
        if v is None: return
        v = float(v)
        seg.update({fmt_br(v, 2), fmt_br(v, 1)})
        lps_v = None
        return lps_v

    def add_speedup(v):
        if v is None: return
        v = float(v)
        x.update({fmt_br(v, 4), fmt_br(v, 3), fmt_br(v, 2)})

    def add_pct(v):
        if v is None: return
        v = float(v) * 100
        pct.update({fmt_br(v, 2), fmt_br(v, 1)})

    def add_lps(v):
        if v is None: return
        v = float(v)
        lps.update({fmt_br(round(v), 0), fmt_br(v, 2)})

    with conn.cursor() as cur:
        cur.execute("SELECT tempo_total_seg, throughput_lps FROM benchmark_execucao")
        for t, thr in cur.fetchall():
            add_tempo(t); add_lps(thr)
        cur.execute("SELECT speedup, eficiencia FROM v_benchmark_metricas WHERE speedup IS NOT NULL")
        for s, e in cur.fetchall():
            add_speedup(s); add_pct(e)
        cur.execute("""SELECT tempo_mediana, tempo_min, tempo_max, tempo_iqr, throughput_mediana,
                              speedup_mediana, speedup_min, speedup_max,
                              eficiencia_mediana, eficiencia_min, eficiencia_max
                       FROM v_benchmark_resumo""")
        for row in cur.fetchall():
            for v in row[:4]: add_tempo(v)
            add_lps(row[4])
            for v in row[5:8]: add_speedup(v)
            for v in row[8:11]: add_pct(v)
        cur.execute("SELECT tempo_seg FROM benchmark_etapa")
        for (t,) in cur.fetchall():
            add_tempo(t)

        # valores citados na EXPLICAÇÃO do bug do baseline (speedup contra o
        # sequencial mais recente) — calculados do banco, como em dados_benchmark
        cur.execute("""SELECT tempo_total_seg FROM benchmark_execucao
                       WHERE campanha_id IS NULL AND modo = 'sequencial' AND NOT aquecimento
                       ORDER BY timestamp_inicio DESC LIMIT 1""")
        _row = cur.fetchone()
        if _row:
            _t_ult = float(_row[0])
            cur.execute("""SELECT tempo_total_seg, num_workers FROM benchmark_execucao
                           WHERE campanha_id IS NULL AND modo = 'paralelo'""")
            for _t_par, _w in cur.fetchall():
                _s_err = _t_ult / float(_t_par)
                add_speedup(_s_err)
                add_pct(_s_err / int(_w))

        # CPU ociosa registrada nas observações das campanhas (condição declarada)
        cur.execute("SELECT DISTINCT observacoes FROM benchmark_execucao "
                    "WHERE observacoes LIKE %s", ("%cpu_ocioso=%",))
        for (obs,) in cur.fetchall():
            for mo in re.finditer(r"cpu_ocioso=([0-9.]+)%", obs or ""):
                _v = float(mo.group(1))
                pct.update({fmt_br(_v, 1), fmt_br(_v, 2)})

    if escalonamento:
        for r in escalonamento.get("agregado", []):
            for k in ("ideal", "teto_real", "teto_lpt", "teto_crescente", "makespan_medido", "wall", "overhead"):
                add_tempo(r.get(k))
            for k in ("speedup_teto_real", "speedup_teto_lpt", "speedup_teto_crescente", "speedup_medido", "inflacao_etapas"):
                add_speedup(r.get(k))
            for k in ("pct_do_teto_real", "pct_do_teto_lpt"):
                add_pct(r.get(k))
            if r.get("cpu_pct") is not None:
                pct.update({fmt_br(float(r["cpu_pct"]), 1), fmt_br(float(r["cpu_pct"]), 2)})
        for l in escalonamento.get("por_execucao", []):
            for k in ("teto_real", "teto_lpt", "teto_crescente", "makespan_medido", "wall", "overhead", "ideal"):
                add_tempo(l.get(k))
            for k in ("speedup_teto_real", "speedup_teto_lpt", "speedup_medido", "inflacao_etapas"):
                add_speedup(l.get(k))
        seqn = escalonamento.get("sequencial", {})
        add_tempo(seqn.get("t_seq_mediana"))
        if seqn.get("cpu_pct_mediana") is not None:
            pct.add(fmt_br(float(seqn["cpu_pct_mediana"]), 1))
        for v in (escalonamento.get("hipoteses") or {}).values():
            if isinstance(v, dict):
                for vv in v.values():
                    if isinstance(vv, (int, float)) and not isinstance(vv, bool):
                        add_speedup(vv)
    return {"s": seg, "x": x, "%": pct, "lps": lps}


def conferir(nome: str, texto: str, permitidos: dict[str, set[str]]) -> tuple[list[str], list[str]]:
    negros: list[str] = []
    nao_rastreaveis: list[str] = []
    achados = {
        "s": RE_SEG.findall(texto),
        "x": RE_X.findall(texto),
        "%": RE_PCT.findall(texto),
        "lps": RE_LPS.findall(texto),
    }
    for classe, tokens in achados.items():
        for tok in tokens:
            if tok in LISTA_NEGRA[classe]:
                negros.append(f"{nome}: {tok} {classe}")
            elif tok not in permitidos[classe]:
                nao_rastreaveis.append(f"{nome}: {tok} {classe}")
    total = sum(len(v) for v in achados.values())
    print(f"{nome}: {total} número(s) de benchmark encontrados "
          f"(s={len(achados['s'])}, ×={len(achados['x'])}, %={len(achados['%'])}, linhas/s={len(achados['lps'])})")
    return negros, sorted(set(nao_rastreaveis))


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--pptx", type=Path, default=None)
    ap.add_argument("--docx", type=Path, default=None)
    ap.add_argument("--escalonamento", type=Path, default=None, help="JSON de analise_escalonamento.py")
    ap.add_argument("--estrito", action="store_true", help="Números não rastreáveis também falham.")
    args = ap.parse_args()
    if not args.pptx and not args.docx:
        ap.error("informe --pptx e/ou --docx")

    esc = json.loads(args.escalonamento.read_text(encoding="utf-8")) if args.escalonamento else None
    conn = conectar()
    try:
        permitidos = valores_permitidos(conn, esc)
    finally:
        conn.close()

    negros_total: list[str] = []
    nr_total: list[str] = []
    if args.pptx:
        n, nr = conferir(args.pptx.name, texto_pptx(args.pptx), permitidos)
        negros_total += n; nr_total += nr
    if args.docx:
        n, nr = conferir(args.docx.name, texto_docx(args.docx), permitidos)
        negros_total += n; nr_total += nr

    if nr_total:
        print(f"\nNúmeros NÃO rastreáveis às views ({len(nr_total)}):")
        for t in nr_total:
            print("  -", t)
    if negros_total:
        print(f"\nNÚMEROS DA LISTA NEGRA PRESENTES ({len(negros_total)}):")
        for t in negros_total:
            print("  -", t)
        return 1
    if args.estrito and nr_total:
        return 1
    print("\nOK: nenhum número da rodada perdida; " +
          ("todos os números rastreados às views." if not nr_total else
           f"{len(nr_total)} não rastreável(is) listados acima (use --estrito para falhar)."))
    return 0


if __name__ == "__main__":
    sys.exit(main())
