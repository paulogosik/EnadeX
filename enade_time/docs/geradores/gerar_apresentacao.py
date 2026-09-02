# -*- coding: utf-8 -*-
"""
Gera a apresentação acadêmica do projeto ENADE-Time Distribuído.

Todos os números de benchmark vêm do banco (views v_benchmark_resumo /
v_benchmark_metricas da campanha oficial) via `dados_benchmark.carregar()`.
Nada é digitado; `verificar_numeros.py` confere o .pptx depois.

Saída: docs/geradores/out/apresentacao.pptx

Uso:
    C:\\Projetos\\ENADE> .venv\\Scripts\\python.exe docs\\geradores\\gerar_apresentacao.py
        [--campanha <uuid>] [--screenshots <pasta>]
"""

import argparse
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from dados_benchmark import MESES, br, carregar, pct, seg, sx  # noqa: E402

AQUI = Path(__file__).resolve().parent
RAIZ = AQUI.parents[1]

ap = argparse.ArgumentParser()
ap.add_argument("--campanha", default="oficial")
ap.add_argument("--screenshots", default=os.environ.get(
    "GERADORES_SCREENSHOTS", str(RAIZ / "apresentacao" / "documento_final" / "screenshots")))
ap.add_argument("--saida", default=str(AQUI / "out" / "apresentacao.pptx"))
args = ap.parse_args()

D = carregar(args.campanha)
SHOTS = Path(args.screenshots)
MAQ = D["maquina"]
ESC = D["escalonamento"]
SEQ = D["sequencial"]
MELHOR = D["melhor"]

# ----------------------------------------------------------------------------
# Paleta (mesma do dashboard) e constantes
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x1D, 0x2C, 0x46)
BLUE = RGBColor(0x2B, 0x44, 0x70)
STEEL = RGBColor(0x44, 0x6E, 0xA3)
TEAL = RGBColor(0x3A, 0x9D, 0x8F)
TEAL_LIGHT = RGBColor(0xD7, 0xEF, 0xEA)
BG = RGBColor(0xF7, 0xF9, 0xFC)
INK = RGBColor(0x22, 0x2B, 0x38)
GRAY = RGBColor(0x6B, 0x77, 0x86)
LIGHT_LINE = RGBColor(0xD8, 0xE0, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SEQ_GRAY = RGBColor(0x94, 0xA3, 0xB8)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)

FONT = "Calibri"
TOTAL_SLIDES = 18
FOOTER_TXT = "ENADE-Time Distribuído — Sistemas Paralelos e Distribuídos"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

_slide_num = 0


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def set_bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def textbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def style_run(run, size, color, bold=False, italic=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic


def add_para(tf, first):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def bullets(slide, items, l, t, w, h, size=17, gap=9, color=INK):
    """items: lista de str ou (nivel, str). 'rotulo:: resto' deixa o rotulo em negrito."""
    tf = textbox(slide, l, t, w, h)
    first = True
    for item in items:
        level, txt = (item if isinstance(item, tuple) else (0, item))
        p = add_para(tf, first)
        first = False
        p.space_after = Pt(gap)
        marker = "•  " if level == 0 else "        –  "
        r = p.add_run()
        r.text = marker
        style_run(r, size, TEAL if level == 0 else GRAY, bold=True)
        if "::" in txt:
            head, rest = txt.split("::", 1)
            r1 = p.add_run()
            r1.text = head
            style_run(r1, size, color, bold=True)
            r2 = p.add_run()
            r2.text = rest
            style_run(r2, size, color)
        else:
            r1 = p.add_run()
            r1.text = txt
            style_run(r1, size, color)
    return tf


def content_slide(title_txt, kicker=None):
    """Cria slide de conteudo com barra de titulo padrao. Retorna slide."""
    global _slide_num
    _slide_num += 1
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.42), Inches(0.12), Inches(0.66)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    tf = textbox(slide, 0.82, 0.30, 11.6, 0.95)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_txt
    style_run(r, 27, NAVY, bold=True)
    if kicker:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = kicker
        style_run(r2, 13, GRAY)

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.22), Inches(12.23), Pt(1.6)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = LIGHT_LINE
    rule.line.fill.background()

    tf_f = textbox(slide, 0.55, 7.06, 9.5, 0.35)
    rf = tf_f.paragraphs[0].add_run()
    rf.text = FOOTER_TXT
    style_run(rf, 9, GRAY)

    tf_n = textbox(slide, 12.0, 7.06, 0.8, 0.35)
    tf_n.paragraphs[0].alignment = PP_ALIGN.RIGHT
    rn = tf_n.paragraphs[0].add_run()
    rn.text = f"{_slide_num} / {TOTAL_SLIDES}"
    style_run(rn, 9, GRAY)
    return slide


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def rbox(slide, l, t, w, h, title, sub=None, fill=BLUE, title_color=WHITE,
         sub_color=None, title_size=14, sub_size=10.5):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    try:
        sp.adjustments[0] = 0.10
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    style_run(r, title_size, title_color, bold=True)
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        style_run(r2, sub_size, sub_color or RGBColor(0xDB, 0xE4, 0xF0))
    return sp


def flow_arrow(slide, l, t, w=0.42, h=0.34, rotation=0):
    ar = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    ar.rotation = rotation
    ar.fill.solid()
    ar.fill.fore_color.rgb = STEEL
    ar.line.fill.background()
    return ar


def make_table(slide, data, l, t, w, h, col_widths=None, font_size=13,
               highlight_row=None):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            elif highlight_row is not None and r == highlight_row:
                cell.fill.fore_color.rgb = TEAL_LIGHT
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 or r == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            bold = (r == 0) or (highlight_row is not None and r == highlight_row)
            style_run(run, font_size, WHITE if r == 0 else INK, bold=bold)
    return table


def picture_if_exists(slide, nome, l, t, w):
    caminho = SHOTS / nome
    if caminho.exists():
        pic = slide.shapes.add_picture(str(caminho), Inches(l), Inches(t), width=Inches(w))
        pic.line.color.rgb = LIGHT_LINE
        pic.line.width = Pt(1)
        return pic
    tf = textbox(slide, l, t, w, 1.0)
    r = tf.paragraphs[0].add_run()
    r.text = f"[captura ausente: {nome}]"
    style_run(r, 11, GRAY, italic=True)
    return None


# Números derivados usados em vários slides ------------------------------------
n_suites = MAQ["n_suites"]
n_completas = MAQ.get("n_suites_completas", n_suites)
suites_rotulo = (f"{n_completas} suítes completas"
                 + (f" (+{n_suites - n_completas} parcial)" if n_suites != n_completas else ""))
w_max = max(D["workers"]) if D["workers"] else None
melhor_esc = D["esc_cfg"](MELHOR["num_workers"], MELHOR["ordem_submissao"]) if MELHOR else None
hip = ESC["hipoteses"]
disco_quente = hip.get("H2_disco_quente_mb_mediana")
disco_aq = D["aquecimento"][0]["disco_mb"] if D["aquecimento"] else None
mes_ano = f"{MESES[D['gerado_em'].month - 1].capitalize()} de {D['gerado_em'].year}"


def h1txt(ordem):
    h = hip.get(f"H1_p6_nao_supera_p4_{ordem}")
    if not h:
        return None
    return (f"{'LPT' if ordem == 'lpt' else 'crescente'}: p=4 {sx(h['speedup_p4'], 2)} × p=6 {sx(h['speedup_p6'], 2)} → "
            f"{'confirmada' if h['confirmada'] else 'refutada'}")


def h3txt(p):
    h = hip.get(f"H3_lpt_ge_crescente_p{p}")
    if not h:
        return None
    return f"p={p}: LPT {sx(h['lpt'], 2)} × crescente {sx(h['crescente'], 2)}"


# ============================================================================
# SLIDE 1 — Capa
# ============================================================================
_slide_num += 1
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.55), Inches(1.5), Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = TEAL
bar.line.fill.background()

tf = textbox(s, 0.9, 1.8, 11.6, 1.6)
r = tf.paragraphs[0].add_run()
r.text = "ENADE-Time Distribuído"
style_run(r, 44, WHITE, bold=True)

tf = textbox(s, 0.9, 2.95, 11.0, 0.9)
r = tf.paragraphs[0].add_run()
r.text = "Sistema Paralelo de Análise Longitudinal dos Microdados do ENADE"
style_run(r, 20, RGBColor(0xB9, 0xC8, 0xDC))

tf = textbox(s, 0.9, 4.45, 11.0, 2.3)
meta = [
    ("Autor(a): ", "Lucas Eduardo Tavares Costa"),
    ("Orientador(a): ", "Fabiano Fagundes"),
    ("Instituição / Curso: ", "CEULP/ULBRA"),
    ("Disciplina: ", "Sistemas Paralelos e Distribuídos"),
    ("", mes_ano),
]
first = True
for head, rest in meta:
    p = add_para(tf, first)
    first = False
    p.space_after = Pt(6)
    if head:
        r1 = p.add_run()
        r1.text = head
        style_run(r1, 14, TEAL, bold=True)
    r2 = p.add_run()
    r2.text = rest
    style_run(r2, 14, RGBColor(0xDB, 0xE4, 0xF0))

notes(s, "Bom dia / boa tarde. Vou apresentar o ENADE-Time Distribuído, projeto da "
         "disciplina de Sistemas Paralelos e Distribuídos. A proposta: medir, em um "
         "caso real com dados públicos do INEP, quanto o paralelismo de dados acelera "
         "o processamento dos microdados do ENADE — e onde está o limite desse ganho, "
         "separando o que é granularidade, o que é escalonamento e o que é contenção.")

# ============================================================================
# SLIDE 2 — Contextualização
# ============================================================================
s = content_slide("Contextualização", "O problema e sua relevância")
bullets(s, [
    "ENADE:: avaliação nacional do INEP; cada edição publica milhões de linhas de microdados.",
    "Processar é custoso:: leitura de TXT em encodings variados, filtragem, validação e carga em banco levam minutos.",
    "Análise longitudinal:: comparar edições (2005–2021) exige consolidar várias bases heterogêneas.",
    "Pergunta de pesquisa:: paralelizar o ETL reduz o tempo total? Quanto? Qual o número ótimo de workers — e por quê?",
    "Por que importa:: caso realista e reprodutível para estudar limites práticos do paralelismo (granularidade, escalonamento, contenção, overhead).",
], 0.85, 1.75, 11.6, 4.6, size=18, gap=14)
notes(s, "Os microdados do ENADE são públicos, volumosos e heterogêneos. O custo de "
         "processar tudo de forma sequencial motiva a pergunta central: vale a pena "
         "paralelizar? O diferencial é medir com rigor — com repetições, máquina "
         "declarada e decomposição das perdas — em vez de assumir que mais workers "
         "significa mais velocidade.")

# ============================================================================
# SLIDE 3 — Objetivos
# ============================================================================
s = content_slide("Objetivos")
tf = textbox(s, 0.85, 1.6, 11.6, 0.6)
r = tf.paragraphs[0].add_run()
r.text = "Objetivo geral"
style_run(r, 18, BLUE, bold=True)
bullets(s, [
    "Aplicar e mensurar paralelismo de dados em um pipeline ETL real sobre os microdados do ENADE.",
], 0.85, 2.1, 11.6, 0.7, size=17)
tf = textbox(s, 0.85, 2.95, 11.6, 0.6)
r = tf.paragraphs[0].add_run()
r.text = "Objetivos específicos"
style_run(r, 18, BLUE, bold=True)
bullets(s, [
    "Construir o pipeline completo: dados brutos → ETL → PostgreSQL → API → dashboard.",
    f"Medir speedup, eficiência e throughput com 1, {', '.join(str(w) for w in D['workers'][:-1])} e {w_max} workers, "
    f"em duas ordens de submissão, com {n_completas} ou mais repetições por configuração.",
    "Separar as perdas de paralelismo em granularidade + escalonamento, contenção e overhead, e testar hipóteses pré-registradas.",
    "Disponibilizar a análise longitudinal em um dashboard interativo (React) sobre API documentada (FastAPI).",
], 0.85, 3.45, 11.6, 3.2, size=17, gap=12)
notes(s, "O objetivo geral não é apenas acelerar, e sim mensurar. Os específicos cobrem "
         "as fases do projeto: infraestrutura de dados, experimento de benchmark com "
         "repetições e hipóteses, e as camadas de exposição (API e dashboard) que "
         "tornam os resultados auditáveis por terceiros.")

# ============================================================================
# SLIDE 4 — Fundamentação
# ============================================================================
s = content_slide("Fundamentação teórica", "Conceitos de paralelismo usados no experimento")
bullets(s, [
    "Paralelismo de dados:: o mesmo código processa partições diferentes do dado (aqui: um ano do ENADE por worker).",
    "Worker:: processo independente do sistema operacional (multiprocessing), evitando o GIL do Python.",
    "Speedup e eficiência:: S(p) = T(1) / T(p);  E(p) = S(p) / p.",
    "Lei de Amdahl:: a fração sequencial f limita o ganho: S(p) ≤ 1 / (f + (1−f)/p).",
    "Escalonamento em lista (Graham, 1969):: com tarefas de tamanhos distintos, a ORDEM de submissão define o makespan; LPT (maior primeiro) é a heurística clássica.",
    "Karp–Flatt:: e = (1/S − 1/p)/(1 − 1/p) estima a fração serial — mas embute granularidade, escalonamento, contenção e overhead.",
], 0.85, 1.7, 11.6, 4.8, size=16.5, gap=12)
notes(s, "Além de Amdahl, dois conceitos que o experimento usa de verdade: o "
         "escalonamento em lista de Graham — o executor entrega cada tarefa ao "
         "primeiro worker livre, na ordem de submissão, então a ordem importa — e a "
         "métrica de Karp–Flatt, que uso com a ressalva de que ela mistura várias "
         "perdas num número só.")

# ============================================================================
# SLIDE 5 — Recorte
# ============================================================================
s = content_slide("Metodologia — recorte do estudo")
bullets(s, [
    "Fonte:: microdados oficiais INEP/ENADE (arquivos TXT, anonimizados — LGPD).",
    "Cursos:: Computação (CO_GRUPO 40 até 2008 e 4004 a partir de 2011 — o mesmo curso, código migrado).",
    "Regiões:: Norte e Nordeste (CO_REGIAO_CURSO 1 e 2).",
    "Edições:: 2005, 2008, 2011, 2014, 2017 e 2021 (ciclos trienais da área).",
    "Base consolidada:: 24.967 inscrições = 582 cursos-ano (a unidade analítica é curso-ano; as linhas replicadas funcionam como peso por matrícula).",
    "Granularidade do paralelismo:: por ano — 6 unidades de trabalho independentes e de tamanhos diferentes.",
], 0.85, 1.7, 11.6, 4.7, size=17, gap=12)
notes(s, "O recorte torna o experimento tratável sem perder representatividade. Dois "
         "pontos honestos: a unidade analítica é curso-ano, 582 observações — as "
         "24.967 linhas são inscrições; e a granularidade por ano dá só 6 unidades de "
         "tamanhos diferentes, o que vai limitar o speedup independentemente de Amdahl.")

# ============================================================================
# SLIDE 6 — Desenho do experimento
# ============================================================================
s = content_slide("Metodologia — desenho do experimento", "Campanha oficial: repetições, duas ordens, máquina e cache declarados")
bullets(s, [
    f"Campanha:: {suites_rotulo}; em cada suíte completa, 1 sequencial + {len(D['workers'])} tamanhos de pool "
    f"({', '.join(str(w) for w in D['workers'])} workers) × 2 ordens de submissão = {MAQ['n_execucoes']} execuções oficiais.",
    "Pareamento por suíte:: cada paralela é comparada com o sequencial da PRÓPRIA suíte (view v_benchmark_metricas); mediana/mín/máx/IQR por configuração em v_benchmark_resumo.",
    "Ordens de submissão:: crescente (2005→2021, ordem histórica) e LPT (maior tempo primeiro, com os tempos do sequencial da suíte).",
    f"Máquina:: {MAQ['cpu_fisicos']} núcleos físicos / {MAQ['cpu_logicos']} lógicos (i5-1135G7), NVMe, 24 GB; "
    f"CPU ociosa antes de medir: {br(MAQ['cpu_ocioso'], 1)} %.",
    f"Cache declarado:: passada de aquecimento descartada (leu {br(disco_aq, 0)} MB do disco); "
    f"execuções oficiais com cache quente ({br(disco_quente, 1)} MB lidos, mediana).",
    "Instrumentação:: CPU média do sistema e bytes lidos do disco por execução (psutil), no processo principal; workers continuam puros.",
], 0.85, 1.7, 11.6, 5.0, size=15.5, gap=11)
notes(s, "Este é o desenho que responde às críticas metodológicas da versão anterior: "
         "repetições em vez de uma medição por ponto; pareamento dentro da suíte em "
         "vez de 'o sequencial mais recente'; as duas ordens de submissão medidas; "
         "máquina com núcleos físicos declarados; e cache quente declarado, não "
         "escondido.")

# ============================================================================
# SLIDE 7 — Arquitetura
# ============================================================================
s = content_slide("Arquitetura", "Fluxo de dados em camadas — pipeline serpentina")
bw, bh = 3.45, 1.25
y1, y2 = 1.9, 4.45
xs = [0.85, 5.0, 9.15]
rbox(s, xs[0], y1, bw, bh, "Microdados INEP", "TXT brutos por edição — somente leitura", fill=STEEL)
flow_arrow(s, 4.42, y1 + 0.45)
rbox(s, xs[1], y1, bw, bh, "ETL Python", "filtro, validação e consolidação (scripts 01–06)", fill=BLUE)
flow_arrow(s, 8.57, y1 + 0.45)
rbox(s, xs[2], y1, bw, bh, "dados_processados/", "CSV consolidado — 24.967 linhas", fill=BLUE)
flow_arrow(s, 10.66, y1 + bh + 0.18, w=0.42, h=0.34, rotation=90)
rbox(s, xs[2], y2, bw, bh, "PostgreSQL 16 (Docker)", "modelo dimensional + benchmark + views de métricas", fill=NAVY)
flow_arrow(s, 8.57, y2 + 0.45, rotation=180)
rbox(s, xs[1], y2, bw, bh, "API FastAPI", "read-only · 19 endpoints · Swagger", fill=BLUE)
flow_arrow(s, 4.42, y2 + 0.45, rotation=180)
rbox(s, xs[0], y2, bw, bh, "Dashboard React", "análises ENADE + métricas SPD", fill=TEAL)
ann = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(6.0), Inches(3.45), Inches(0.78))
try:
    ann.adjustments[0] = 0.12
except Exception:
    pass
ann.fill.solid()
ann.fill.fore_color.rgb = TEAL_LIGHT
ann.line.color.rgb = TEAL
ann.line.width = Pt(1)
tf = ann.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Benchmark (scripts 08–10): mede tempo por etapa e grava no próprio banco"
style_run(r, 10.5, RGBColor(0x14, 0x55, 0x4B), bold=True)
notes(s, "Fluxo em serpentina: brutos → ETL → CSV → Postgres via COPY → API → dashboard. "
         "O benchmark instrumenta o ETL e grava no próprio banco; as views calculam as "
         "métricas em SQL, então dashboard e documento leem a mesma definição.")

# ============================================================================
# SLIDE 8 — Tecnologias
# ============================================================================
s = content_slide("Tecnologias e decisões de projeto")
make_table(s, [
    ["Camada", "Tecnologia", "Decisão-chave"],
    ["ETL / Benchmark", "Python 3.13 · pandas · multiprocessing · psycopg2 · psutil", "processos (não threads) por causa do GIL; workers puros"],
    ["Banco de dados", "PostgreSQL 16 em Docker Compose", "COPY FROM STDIN; métricas em views SQL auditáveis"],
    ["API", "FastAPI · Pydantic v2 · pool de conexões", "read-only; sem ORM — a API só lê as views"],
    ["Frontend", "React 18 · TypeScript · Vite · TanStack Query", "tipos TS espelham os schemas Pydantic"],
    ["Visualização", "Recharts · Tailwind CSS", "mediana ± mín–máx; filtros deep-link na URL"],
], 0.85, 1.75, 11.65, 3.9, col_widths=[2.3, 4.85, 4.5], font_size=13)
notes(s, "Sem ORM na API — SQL explícito com placeholders. A definição de speedup vive "
         "em uma única view; a API não recalcula nada, só lê. No frontend, as "
         "interfaces TypeScript espelham os schemas Pydantic.")

# ============================================================================
# SLIDE 9 — Fases
# ============================================================================
s = content_slide("Desenvolvimento — fases do projeto")
bullets(s, [
    "Fase 1:: schema dimensional no PostgreSQL + carga do CSV consolidado (COPY) — 24.967 registros validados.",
    "Fase 2:: experimento de benchmark — sequencial × paralelo, com instrumentação por etapa.",
    "Fase 3:: API FastAPI read-only — 19 endpoints (dimensões, análises, benchmark, campanhas).",
    "Fase 4:: dashboard React — 11 páginas, KPIs em tempo real, comparativo SPD e drill-down por worker.",
    "Fase 5:: documentação — guia de execução, arquitetura, resultados e roteiro de apresentação.",
    "Fase 6 (v2):: correção do baseline do comparativo, campanha oficial com repetições e duas ordens, integração ao ecossistema EnadeX.",
], 0.85, 1.7, 11.6, 4.7, size=16.5, gap=13)
notes(s, "Construção incremental. A fase 6 nasceu de um problema real: o comparativo do "
         "dashboard dividia toda execução pelo sequencial mais recente, e ao misturar "
         "duas rodadas produziu um speedup inflado. Corrigir isso exigiu repensar o "
         "experimento — e o resultado é a campanha que vem a seguir.")

# ============================================================================
# SLIDE 10 — Resultados (tabela + gráfico)
# ============================================================================
s = content_slide("Resultados — campanha oficial",
                  f"{suites_rotulo} · mediana [mín–máx] · pareamento por suíte (n por configuração na tabela)")

linhas = [["Configuração", "n", "Tempo (s)", "Speedup", "Eficiência", "Throughput (linhas/s)"]]
melhor_idx = None
for i, r_ in enumerate(D["resumo"], start=1):
    linhas.append([
        r_["rotulo"], str(r_["n"]),
        f"{br(r_['tempo_mediana'])} [{br(r_['tempo_min'])}–{br(r_['tempo_max'])}]",
        sx(r_["speedup_mediana"], 3) if r_["modo"] == "paralelo" else "1,000×",
        pct(r_["eficiencia_mediana"]) if r_["modo"] == "paralelo" else "100,0 %",
        br(r_["throughput_mediana"], 0),
    ])
    if MELHOR and r_ is MELHOR:
        melhor_idx = i
make_table(s, linhas, 0.85, 1.55, 7.2, 3.6,
           col_widths=[2.0, 0.45, 1.85, 0.95, 0.95, 1.0], font_size=10.5, highlight_row=melhor_idx)

chart_data = CategoryChartData()
chart_data.categories = [("Seq." if r_["modo"] == "sequencial" else
                          f"{r_['num_workers']}w {'LPT' if r_['ordem_submissao'] == 'lpt' else 'cresc.'}")
                         for r_ in D["resumo"]]
chart_data.add_series("Tempo mediano (s)", tuple(round(r_["tempo_mediana"], 2) for r_ in D["resumo"]))
gframe = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.25), Inches(1.55), Inches(4.5), Inches(3.6), chart_data
)
chart = gframe.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 60
plot.has_data_labels = True
dl = plot.data_labels
dl.number_format = "0.0"
dl.number_format_is_linked = False
dl.font.size = Pt(9)
dl.font.bold = True
dl.font.color.rgb = INK
series = plot.series[0]
for i, r_ in enumerate(D["resumo"]):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = (SEQ_GRAY if r_["modo"] == "sequencial"
                                     else (TEAL if r_["ordem_submissao"] == "lpt" else VIOLET))
chart.category_axis.tick_labels.font.size = Pt(9)
chart.value_axis.tick_labels.font.size = Pt(9)

card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.35), Inches(11.9), Inches(1.55))
try:
    card.adjustments[0] = 0.08
except Exception:
    pass
card.fill.solid()
card.fill.fore_color.rgb = WHITE
card.line.color.rgb = TEAL
card.line.width = Pt(1.5)
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.22)
tf.margin_right = Inches(0.22)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
r = p.add_run()
if MELHOR:
    r.text = (f"Melhor configuração: {MELHOR['rotulo']} — speedup mediano {sx(MELHOR['speedup_mediana'])} "
              f"[{sx(MELHOR['speedup_min'], 2)}–{sx(MELHOR['speedup_max'], 2)}], eficiência {pct(MELHOR['eficiencia_mediana'])}")
else:
    r.text = "Sem execuções paralelas na campanha"
style_run(r, 15, RGBColor(0x14, 0x55, 0x4B), bold=True)
for line in [
    f"Baseline sequencial: mediana {seg(SEQ['tempo_mediana'])} [{br(SEQ['tempo_min'])}–{br(SEQ['tempo_max'])}] em {SEQ['n']} medições; "
    f"IQR {br(SEQ['tempo_iqr'])} s.",
    (f"Throughput máximo: {br(MELHOR['throughput_mediana'], 0)} linhas/s (mediana) × {br(SEQ['throughput_mediana'], 0)} no sequencial."
     if MELHOR else ""),
]:
    if not line:
        continue
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r1 = p2.add_run()
    r1.text = "•  "
    style_run(r1, 12, TEAL, bold=True)
    r2 = p2.add_run()
    r2.text = line
    style_run(r2, 12, INK)

notes(s, f"Números da campanha oficial (id {D['campanha_id'][:8]}), {suites_rotulo}. Cada linha é a "
         f"mediana das repetições (n na tabela), com mínimo e máximo. O sequencial mediano é "
         f"{seg(SEQ['tempo_mediana'])}. A melhor configuração foi {MELHOR['rotulo'] if MELHOR else '—'}, "
         f"com speedup mediano {sx(MELHOR['speedup_mediana']) if MELHOR else '—'}. Tudo isto está em "
         f"v_benchmark_resumo — quem quiser confere com um SELECT.")

# ============================================================================
# SLIDE 11 — Análise: três perdas
# ============================================================================
s = content_slide("Análise — onde o speedup se perde", "Tetos por granularidade + escalonamento, contenção e overhead (medidos)")

tetos = [["p", "Teto ordem crescente", "Teto LPT", "Medido cresc.", "Medido LPT", "% do teto (LPT)", "Inflação etapas", "Overhead (s)"]]
for w in D["workers"]:
    ec, el = D["esc_cfg"](w, "crescente"), D["esc_cfg"](w, "lpt")
    rc, rl = D["cfg"](w, "crescente"), D["cfg"](w, "lpt")
    tetos.append([
        str(w),
        sx((ec or el or {}).get("speedup_teto_crescente"), 2),
        sx((el or ec or {}).get("speedup_teto_lpt"), 2),
        sx(rc["speedup_mediana"], 2) if rc else "—",
        sx(rl["speedup_mediana"], 2) if rl else "—",
        pct(el["pct_do_teto_lpt"], 0) if el and el.get("pct_do_teto_lpt") is not None else "—",
        f"{br(el['inflacao_etapas'], 2)}×" if el and el.get("inflacao_etapas") is not None else "—",
        br(el["overhead"], 1) if el and el.get("overhead") is not None else "—",
    ])
make_table(s, tetos, 0.85, 1.55, 11.9, 2.1,
           col_widths=[0.5, 1.75, 1.35, 1.35, 1.3, 1.45, 1.4, 1.3], font_size=11)

analise_itens = [
    "Granularidade + escalonamento:: 6 unidades de tamanhos diferentes; o executor entrega cada ano ao primeiro worker livre, "
    "então a ordem de submissão fixa o teto (Graham). Acima de 6 workers o teto não cresce.",
    (f"Contenção:: com {MAQ['cpu_fisicos']} núcleos físicos, cada etapa fica ~{br(melhor_esc['inflacao_etapas'], 2)}× mais lenta "
     f"em paralelo ({MELHOR['rotulo']}) — a perda não é de disco: {br(disco_quente, 1)} MB lidos por execução quente "
     f"contra {br(disco_aq, 0)} MB no aquecimento."
     if melhor_esc and melhor_esc.get("inflacao_etapas") is not None else
     "Contenção:: medida pela inflação das etapas em paralelo (soma dos tempos por ano ÷ soma no sequencial)."),
    (f"Overhead:: spawn dos processos, coleta e gravação ≈ {br(melhor_esc['overhead'], 1)} s por execução — pequeno frente ao makespan."
     if melhor_esc and melhor_esc.get("overhead") is not None else "Overhead:: wall-clock − makespan medido."),
]
h1s = [t for t in (h1txt("lpt"), h1txt("crescente")) if t]
if h1s:
    analise_itens.append("H1 (p=6 não supera p=4 — hyperthreads):: " + "; ".join(h1s) + ".")
analise_itens.append(f"H2 (cache quente → limitado por CPU, não por disco):: {br(disco_quente, 1)} MB lidos nas execuções quentes; "
                     f"CPU média {br(MAQ['cpu_percent_medio'], 1)} % do sistema.")
h3s = [t for t in (h3txt(2), h3txt(3), h3txt(4)) if t]
if h3s:
    analise_itens.append("H3 (LPT ≥ crescente):: " + "; ".join(h3s) + ".")
bullets(s, analise_itens, 0.85, 3.85, 11.9, 3.1, size=13, gap=7)
notes(s, "Este slide substitui a explicação antiga ('4 workers lendo do mesmo disco geram "
         "contenção'), que era suposição. Agora: a parte do teto é matemática (tempos "
         "por ano + ordem de submissão); a contenção aparece como inflação das etapas "
         "com o disco praticamente parado — é CPU em 4 núcleos físicos; e o overhead "
         "de spawn é pequeno. Karp–Flatt embutia as três coisas num número só.")

# ============================================================================
# SLIDE 12 — Histórico das medições e o bug do baseline
# ============================================================================
s = content_slide("Histórico das medições — e o erro que a campanha corrige",
                  "Três rodadas anteriores; uma delas não existe mais em banco; uma foi lida com o baseline errado")
hist_rows = [["Rodada", "Sequencial (s)", "2 workers", "4 workers", "Situação"]]
for rod in D["historico"]:
    b = rod["baseline"]
    p2 = next((x for x in rod["paralelas"] if x["num_workers"] == 2), None)
    p4 = next((x for x in rod["paralelas"] if x["num_workers"] == 4), None)
    hist_rows.append([
        b["timestamp_inicio"].strftime("%d/%m/%Y") + f" (#{b['id']})",
        br(b["tempo_total_seg"]),
        f"{br(p2['tempo_total_seg'])} s · {sx(p2['speedup'], 4)}" if p2 else "—",
        f"{br(p4['tempo_total_seg'])} s · {sx(p4['speedup'], 4)}" if p4 else "—",
        "no banco; pareamento temporal (ordem crescente)",
    ])
hist_rows.append(["Apresentação anterior", "—", "—", "—",
                  "rodada apagada num --reset do schema: sem suporte em banco; não é mais citada"])
make_table(s, hist_rows, 0.85, 1.55, 11.9, 0.45 * len(hist_rows),
           col_widths=[2.1, 1.5, 2.6, 2.6, 3.1], font_size=11)

bug_itens = []
for b in D["bug_baseline"]:
    bug_itens.append(
        f"Execução #{b['exec']} ({b['workers']} workers, {seg(b['tempo'])}):: o comparativo antigo dividia pelo sequencial "
        f"mais recente (#{b['baseline_errado']}, {seg(b['t_errado'])}) → {sx(b['s_errado'])} / {pct(b['e_errado'])}; "
        f"pareada com o próprio baseline (#{b['baseline_certo']}) → {sx(b['s_certo'])} / {pct(b['e_certo'])}.")
bug_itens.append("Correção:: pareamento por suíte na view; a API só lê a view; checagem cruzada em Python "
                 "(scripts/13_validar_metricas.py) e nos testes; números de documento e slides conferidos "
                 "contra as views por script.")
bug_itens.append("Lição:: a medição importa tanto quanto o código — duas vezes (OneDrive em 2026-06 e baseline em 2026-08).")
bullets(s, bug_itens, 0.85, 1.7 + 0.45 * len(hist_rows) + 0.15, 11.9, 3.2, size=13, gap=8)
notes(s, "Transparência sobre o histórico: a rodada da apresentação anterior foi apagada "
         "por um reset do schema e não tem suporte em banco; as duas rodadas que "
         "existem foram misturadas pelo comparativo, que usava o sequencial mais "
         "recente como baseline de tudo — daí um speedup inflado. A correção está na "
         "view, na API e nos testes.")

# ============================================================================
# SLIDE 13 — Dashboard
# ============================================================================
s = content_slide("Resultados — dashboard interativo")
bullets(s, [
    "Visão geral:: KPIs reais via API (registros, edições, baseline mediano, melhor speedup).",
    "Análises ENADE:: média por ano, região, UF e ranking de IES — com filtros combináveis.",
    "Comparativo SPD:: mediana ± mín–máx por configuração; uma linha de eficiência por ordem de submissão.",
    "Condições de medição na tela:: núcleos físicos/lógicos, cache quente, campanha e nº de suítes.",
    "Drill-down:: tempo por etapa/ano e por worker de cada execução; histórico com pareamento explícito.",
    "API documentada:: Swagger em /docs; /benchmark/campanhas e /comparativo?campanha_id=… para auditoria.",
], 0.85, 1.7, 6.0, 5.0, size=15, gap=10)
picture_if_exists(s, "12_spd_comparativo.png", 7.15, 1.75, 5.4)
notes(s, "Demo ao vivo ou captura. O texto de interpretação é gerado dos dados do banco — "
         "se a campanha mudar, o texto muda. A execução ao vivo da demo entra como não "
         "oficial, e a diferença para a mediana oficial é o tamanho do que o ambiente "
         "de apresentação contamina.")

# ============================================================================
# SLIDE 14 — Telas
# ============================================================================
s = content_slide("Dashboard — telas principais", "Capturas do sistema em execução")
_telas = [
    ("09_home.png", "Visão geral — KPIs via API"),
    ("11_comparar.png", "Comparação Norte × Nordeste"),
    ("13_etapas.png", "Drill-down por worker"),
]
_x = 0.55
for _arq, _rotulo in _telas:
    _pic = picture_if_exists(s, _arq, _x, 1.9, 4.0)
    alt = 1.9 + (4.0 * (_pic.height / _pic.width) if _pic else 2.5) + 0.12
    _tf = textbox(s, _x, alt, 4.0, 0.5)
    _p = _tf.paragraphs[0]
    _p.alignment = PP_ALIGN.CENTER
    _r = _p.add_run()
    _r.text = _rotulo
    style_run(_r, 12, GRAY, bold=True)
    _x += 4.25
notes(s, "Três telas: visão geral; comparação longitudinal Norte × Nordeste com exportação; "
         "drill-down mostrando qual worker processou cada ano — é aqui que se vê a ordem "
         "de submissão em ação.")

# ============================================================================
# SLIDE 15 — Conclusão
# ============================================================================
s = content_slide("Conclusão")
bullets(s, [
    "Pipeline completo e reprodutível:: dos TXT do INEP ao dashboard — qualquer revisor reexecuta do zero.",
    (f"Ganho real, com teto conhecido:: speedup mediano {sx(MELHOR['speedup_mediana'], 2)} com {MELHOR['rotulo']} "
     f"({pct(melhor_esc['pct_do_teto_real'], 0) if melhor_esc and melhor_esc.get('pct_do_teto_real') is not None else '—'} do teto da ordem usada)."
     if MELHOR else "Ganho real, com teto conhecido."),
    "Três perdas, três causas:: granularidade + escalonamento (matemática dos 6 anos), contenção em 4 núcleos físicos (disco quase parado), overhead pequeno.",
    "Medir > supor, três vezes:: OneDrive contaminou o I/O; o baseline errado inflou um speedup; a explicação por 'disco' não sobreviveu aos bytes lidos.",
    "Arquitetura em camadas:: definição única das métricas em SQL; API, dashboard, documento e slides leem a mesma view — e um script confere.",
], 0.85, 1.7, 11.6, 4.8, size=16, gap=13)
notes(s, "Mensagens para levar: paralelismo tem teto e tem custo; a metodologia de "
         "medição importa tanto quanto o código; e uma arquitetura em camadas com uma "
         "única definição das métricas torna o experimento auditável.")

# ============================================================================
# SLIDE 16 — Limitações e trabalhos futuros
# ============================================================================
s = content_slide("Limitações e trabalhos futuros")
tf = textbox(s, 0.85, 1.6, 5.7, 0.5)
r = tf.paragraphs[0].add_run()
r.text = "Limitações"
style_run(r, 18, BLUE, bold=True)
bullets(s, [
    f"Uma única máquina (Windows, {MAQ['cpu_fisicos']} núcleos físicos); CPU ociosa de {br(MAQ['cpu_ocioso'], 1)} % antes de medir.",
    "Granularidade fixa (por ano) — 6 unidades; p = 8 excluído: dois workers ociosos por construção.",
    "Cache de páginas quente em todas as execuções oficiais (declarado; a versão fria é o aquecimento).",
    f"{n_completas}+ repetições por configuração: mín–máx e IQR, não intervalo de confiança formal.",
], 0.85, 2.15, 5.8, 3.8, size=14.5, gap=9)
tf = textbox(s, 7.0, 1.6, 5.7, 0.5)
r = tf.paragraphs[0].add_run()
r.text = "Trabalhos futuros"
style_run(r, 18, BLUE, bold=True)
bullets(s, [
    "Granularidade menor (partições dentro do ano) para elevar o teto de escalonamento.",
    "Repetir a campanha em máquina com mais núcleos físicos e em Linux (fork vs spawn).",
    "Eixo de carga: COPY local × COPY via pooler × upsert PostgREST (custo de publicar na nuvem).",
    "Profiling por etapa (py-spy) para separar parsing de I/O dentro de cada ano.",
    "Docker completo, CI com a suíte de testes e integração final no ecossistema EnadeX.",
], 7.0, 2.15, 5.8, 4.2, size=14.5, gap=9)
notes(s, "Limitações reconhecidas: uma máquina, cache quente, poucas repetições para "
         "um IC formal. O próximo passo mais promissor é reduzir a granularidade — é o "
         "teto de escalonamento, não Amdahl, que segura o speedup aqui.")

# ============================================================================
# SLIDE 17 — Referências
# ============================================================================
s = content_slide("Referências")
bullets(s, [
    "AMDAHL, G. M. Validity of the single processor approach to achieving large scale computing capabilities. AFIPS, 1967.",
    "GRAHAM, R. L. Bounds on multiprocessing timing anomalies. SIAM Journal on Applied Mathematics, 17(2), 1969.",
    "KARP, A. H.; FLATT, H. P. Measuring parallel processor performance. Communications of the ACM, 33(5), 1990.",
    "PACHECO, P. An Introduction to Parallel Programming. Morgan Kaufmann, 2011.",
    "INEP. Microdados do ENADE — edições 2005 a 2021. Brasília: INEP/MEC. Disponível em: gov.br/inep.",
    "PostgreSQL 16 Documentation. postgresql.org/docs/16. · FastAPI Documentation. fastapi.tiangolo.com. · React Documentation. react.dev.",
    "Documentação interna: docs/ARQUITETURA.md, docs/RESULTADOS_BENCHMARK.md, docs/RELATORIO_VERIFICACAO.md, docs/DESIGN_LOG.md.",
], 0.85, 1.75, 11.6, 4.8, size=13.5, gap=10)
notes(s, "Referências teóricas (Amdahl, Graham, Karp e Flatt, Pacheco), a fonte oficial "
         "dos dados e a documentação interna do projeto.")

# ============================================================================
# SLIDE 18 — Encerramento
# ============================================================================
_slide_num += 1
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.0), Inches(1.5), Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = TEAL
bar.line.fill.background()
tf = textbox(s, 0.9, 2.25, 11.0, 1.0)
r = tf.paragraphs[0].add_run()
r.text = "Obrigado!"
style_run(r, 44, WHITE, bold=True)
tf = textbox(s, 0.9, 3.35, 11.0, 0.6)
r = tf.paragraphs[0].add_run()
r.text = "Perguntas?"
style_run(r, 22, RGBColor(0xB9, 0xC8, 0xDC))
tf = textbox(s, 0.9, 4.7, 11.5, 2.2)
info = [
    ("Repositório: ", "github.com/paulogosik/EnadeX — pasta enade_time/"),
    ("Dashboard: ", "http://localhost:5173"),
    ("API / Swagger: ", "http://localhost:8002/docs"),
    ("Campanha oficial: ", D["campanha_id"]),
]
first = True
for head, rest in info:
    p = add_para(tf, first)
    first = False
    p.space_after = Pt(6)
    r1 = p.add_run()
    r1.text = head
    style_run(r1, 14, TEAL, bold=True)
    r2 = p.add_run()
    r2.text = rest
    style_run(r2, 14, RGBColor(0xDB, 0xE4, 0xF0))
notes(s, "Agradecer e abrir para perguntas. Respostas prováveis: com mais núcleos físicos o "
         "ponto ótimo desloca, mas o teto de escalonamento dos 6 anos permanece; a unidade "
         "analítica é curso-ano (582); os números são auditáveis pela view SQL.")

# ----------------------------------------------------------------------------
# Salvar e validar
# ----------------------------------------------------------------------------
prs.core_properties.title = "ENADE-Time Distribuído — Apresentação acadêmica"
prs.core_properties.subject = "Sistema Paralelo de Análise Longitudinal dos Microdados do ENADE"
prs.core_properties.author = "Lucas Eduardo Tavares Costa"

out = Path(args.saida)
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
check = Presentation(str(out))
n_slides = len(list(check.slides))
print("OK: arquivo gerado em", out)
print(f"Slides: {n_slides} (esperado: {TOTAL_SLIDES}); campanha {D['campanha_id']}")
if n_slides != TOTAL_SLIDES:
    sys.exit(1)
